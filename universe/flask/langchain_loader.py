import nltk
from langchain.schema import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import Chroma
import os
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document as WordDocument
# Set nltk data path and download dependencies
nltk.data.path.append('/home/amogh/nltk_data')
nltk.download('punkt')

CHROMA_PATH = "chroma"
DATA_PATH = "uploads"
PROCESSED_FILES_TRACKER = "processed_files.txt"  # File to keep track of processed files


def generate_data_store(username):
    try:
        documents = load_new_documents(username)
        if documents:
            chunks = split_text(documents)
            save_to_chroma(chunks)
        else:
            print("No new documents to process.")
    except Exception as e:
        print(f"Error in generate_data_store: {e}")


def load_existing_files():
    try:
        if os.path.exists(PROCESSED_FILES_TRACKER):
            with open(PROCESSED_FILES_TRACKER, "r") as f:
                existing_files = set(f.read().splitlines())
        else:
            existing_files = set()
        return existing_files
    except Exception as e:
        print(f"Error in load_existing_files: {e}")
        return set()


def update_processed_files(new_files):
    try:
        with open(PROCESSED_FILES_TRACKER, "a") as f:
            f.write("\n".join(new_files) + "\n")
    except Exception as e:
        print(f"Error in update_processed_files: {e}")


def load_new_documents(username):
    try:
        existing_files = load_existing_files()
        all_files = {f for f in os.listdir(DATA_PATH+f"_{username}") if f.endswith(('.pdf', '.ppt', '.pptx', '.docx'))}
        new_files = all_files - existing_files

        documents = []
        for file_name in new_files:
            file_path = os.path.join(DATA_PATH+f"_{username}", file_name)
            text = ""

            if file_name.endswith('.pdf'):
                # Extract text from PDF
                with open(file_path, 'rb') as f:
                    reader = PdfReader(f)
                    text = "".join([page.extract_text() for page in reader.pages])

            elif file_name.endswith(('.ppt', '.pptx')):
                # Extract text from PowerPoint
                presentation = Presentation(file_path)
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            text += "\n".join([p.text for p in shape.text_frame.paragraphs])

            elif file_name.endswith('.docx'):
                # Extract text from Word document
                doc = WordDocument(file_path)
                text = "\n".join([paragraph.text for paragraph in doc.paragraphs])

            if text.strip():  # Only add non-empty documents
                documents.append(Document(page_content=text, metadata={"source": file_name}))

        if documents:
            update_processed_files(new_files)

        return documents
    except Exception as e:
        print(f"Error in load_new_documents: {e}")
        return []

def split_text(documents: list[Document]):
    try:
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=300,
            chunk_overlap=100,
            length_function=len,
            add_start_index=True
        )
        chunks = text_splitter.split_documents(documents)
        print(f"Split {len(documents)} documents into {len(chunks)} chunks.")

        # Example chunk preview
        if chunks:
            document = chunks[0]  # Preview the first chunk
            print("Example Chunk Content:")
            print(document.page_content)
            print("Metadata:")
            print(document.metadata)

        return chunks
    except Exception as e:
        print(f"Error in split_text: {e}")
        return []


def save_to_chroma(chunks: list[Document]):
    try:
        embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L12-v2")

        if os.path.exists(CHROMA_PATH):
            db = Chroma(persist_directory=CHROMA_PATH, embedding_function=embeddings)
        else:
            db = Chroma.from_documents([], embeddings, persist_directory=CHROMA_PATH)

        # Add new chunks to the database
        db.add_documents(chunks)
        db.persist()
        print(f"Saved {len(chunks)} new chunks to {CHROMA_PATH}")
    except Exception as e:
        print(f"Error in save_to_chroma: {e}")